#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ABB 로봇 시스템 기술 검토 - PowerPoint 생성 스크립트 (한글)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """PowerPoint 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 1. 타이틀 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (25, 25, 112), (65, 105, 225))

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "ABB 로봇 시스템\n기술 검토"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Controller: 1200-526402 | RobotWare 6.16"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(200, 200, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 2. 시스템 개요
    slide = add_title_content_slide(prs, "시스템 개요")
    add_bullet_points(slide, [
        "프로젝트: S25016 34bay 자동용접 라인",
        "제어기: IRC5 (1200-526402)",
        "RobotWare: 6.16.01",
        "구성: 3개 로봇 협조 제어",
        "  • Robot1, Robot2: IRB 1200-5/0.9 (용접 로봇)",
        "  • Robot3 (Gantry): 4축 외부축 (X, Y, Z, R)",
        "특징: MultiMove 기반 준동기화 시스템"
    ])

    # 3. 로봇 구성
    slide = add_title_content_slide(prs, "로봇 구성")
    add_bullet_points(slide, [
        "Robot1 (IRB 1200-5/0.9)",
        "  • Base Frame: [4665, -2217, 0]",
        "  • 용접 툴: tWeld1",
        "  • RAPID Task: T_ROB1",
        "",
        "Robot2 (IRB 1200-5/0.9)",
        "  • Base Frame: [4665, 2217, 0]",
        "  • 용접 툴: tWeld2",
        "  • RAPID Task: T_ROB2",
        "",
        "Robot3 (Gantry - 외부축)",
        "  • 4축 갠트리: X, Y, Z 이동 + R 회전",
        "  • RAPID Task: Gantry_T"
    ])

    # 4. 제어기 정보
    slide = add_title_content_slide(prs, "제어기 정보")
    add_bullet_points(slide, [
        "Controller Type: IRC5 Single Cabinet",
        "Serial Number: 1200-526402",
        "RobotWare Version: 6.16.01",
        "",
        "주요 옵션:",
        "  • MultiMove Coordinated (3개 메카니컬 유닛)",
        "  • Arc Welding (용접 전용)",
        "  • Collision Detection",
        "  • Path Recovery",
        "  • World Zones",
        "",
        "DeviceNet: 구성되어 있으나 실제 사용 안 함",
        "  → Lincoln 용접기는 ArcLink-XT 사용"
    ])

    # 5. RAPID 프로그램 구조
    slide = add_title_content_slide(prs, "RAPID 프로그램 구조")
    add_bullet_points(slide, [
        "총 10개 Task 구성:",
        "",
        "Task 1 (T_ROB1) - Robot1 용접",
        "Task 2 (T_ROB2) - Robot2 용접",
        "Task 7 (Gantry_T) - 갠트리 제어",
        "Task 3 (T_ROB1_SCAN) - Robot1 스캐닝",
        "Task 4 (T_ROB2_SCAN) - Robot2 스캐닝",
        "Task 5-6, 8-10 - 예비 (코드 없음)",
        "",
        "총 라인 수: 약 5,000+ 라인",
        "핵심 모듈: MainModule, WeldModule, SafetyModule"
    ])

    # 6. 핵심 기술 1: TCP 좌표 변환
    slide = add_title_content_slide(prs, "핵심 기술 1: TCP 좌표 변환")
    add_two_column_content(slide,
        left_title="문제 상황",
        left_content=[
            "ABB MultiMove는 최대 2축까지만",
            "하드웨어 동기화 지원",
            "",
            "3개 로봇(Robot1 + Robot2 + Gantry)",
            "동시 제어 필요",
            "",
            "→ 소프트웨어 좌표 변환 필요"
        ],
        right_title="해결 방법: RAPID 함수",
        right_content=[
            "fnPoseToExtax(robtarget)",
            "  • TCP 좌표 → 외부축 좌표",
            "  • X, Y, Z, R(회전) 추출",
            "  • 4축 모두 변환",
            "",
            "fnCoordToJoint(jointtarget)",
            "  • 홈 위치 오프셋 적용",
            "  • 안전 리미트 체크",
            "  • 부호 변환 (Y, Z, R 축)"
        ])

    # 7. 핵심 기술 2: 3-Task 준동기화
    slide = add_title_content_slide(prs, "핵심 기술 2: 3-Task 준동기화")
    add_two_column_content(slide,
        left_title="동기화 메커니즘",
        left_content=[
            "WaitSyncTask",
            "  • 3개 Task를 동일 시점 대기",
            "  • taskGroup123 정의",
            "",
            "nWeldSequence 변수",
            "  • 1 → 2 → 3 단계 제어",
            "  • Task 간 상태 공유",
            "",
            "WaitUntil 조건 대기",
            "  • nWeldSequence=3 확인",
            "  • Gantry 이동 시작 타이밍"
        ],
        right_title="실행 흐름",
        right_content=[
            "1. 모든 Task WaitSyncTask 도달",
            "",
            "2. nWeldSequence = 1 설정",
            "   (Robot 진입 단계)",
            "",
            "3. nWeldSequence = 2 변경",
            "   (Gantry 이동 시작)",
            "",
            "4. nWeldSequence = 3 변경",
            "   (Robot 마무리)",
            "",
            "→ 하드웨어 동기화 없이",
            "   소프트웨어로 협조 제어"
        ])

    # 8. 핵심 기술 3: 3단계 모션 시퀀스
    slide = add_title_content_slide(prs, "핵심 기술 3: 3단계 모션 시퀀스")
    add_bullet_points(slide, [
        "Phase 1: 로봇 진입 (nWeldSequence = 1)",
        "  • Robot1, Robot2가 용접 시작점으로 이동",
        "  • 아크 점화 및 초기 용접",
        "  • Gantry는 정지 상태 유지",
        "",
        "Phase 2: 갠트리 이동 (nWeldSequence = 2)",
        "  • Gantry가 Y축 방향으로 이동 시작",
        "  • Robot1, Robot2는 매우 느린 속도로 위빙 유지",
        "  • so_MoveG_PosHold = 1 신호 활성화",
        "",
        "Phase 3: 로봇 마무리 (nWeldSequence = 3)",
        "  • Gantry 이동 완료 및 정지",
        "  • Robot1, Robot2가 용접 마무리",
        "  • 아크 소호 및 복귀"
    ])

    # 9. 핵심 기술 4: 위빙 제어
    slide = add_title_content_slide(prs, "핵심 기술 4: 갠트리 이동 중 위빙")
    add_two_column_content(slide,
        left_title="기술적 과제",
        left_content=[
            "문제:",
            "  • Gantry가 Robot을 Y축으로 이동",
            "  • 동시에 Robot은 X축 위빙 필요",
            "  • 대각선 용접선 생성",
            "",
            "제약:",
            "  • 하드웨어 동기화 불가",
            "  • 좌표계 충돌 위험",
            "  • 속도 매칭 필요"
        ],
        right_title="해결 방법",
        right_content=[
            "so_MoveG_PosHold 신호",
            "  • Gantry 이동 상태 표시",
            "  • Robot에게 전달",
            "",
            "WHILE 루프 제어",
            "  • 신호=1 동안 위빙 유지",
            "  • 매우 느린 속도 (vWeld{4})",
            "  • Gantry가 실제 Y축 이동 담당",
            "",
            "결과:",
            "  • X축 위빙 + Y축 이동",
            "  • 안정적인 대각선 용접"
        ])

    # 10. DeviceNet 구성
    slide = add_title_content_slide(prs, "DeviceNet 구성")
    add_bullet_points(slide, [
        "현재 상태: 구성되어 있으나 사용 안 함",
        "",
        "설정 정보:",
        "  • Device 10: PowerWave-R450 (Robot1)",
        "  • Device 11: PowerWave-R450 (Robot2)",
        "  • Polling 158 bytes, UCMM 512 bytes",
        "",
        "실제 통신 방식:",
        "  • Lincoln ArcLink-XT 소프트웨어 인터페이스 사용",
        "  • DeviceNet 하드웨어 사용 안 함",
        "",
        "향후 계획:",
        "  • DeviceNet 설정 제거 검토",
        "  • 또는 실제 활성화 고려"
    ])

    # 11. 용접 파라미터
    slide = add_title_content_slide(prs, "용접 모션 파라미터")
    add_two_column_content(slide,
        left_title="용접 조건 (welddata)",
        left_content=[
            "속도별 용접 조건 (40개 세트):",
            "",
            "wd1~wd4: 8.66 mm/s",
            "  • 전압: 40V",
            "  • 전류: 270A",
            "  • 와이어 송급: 345 ipm",
            "",
            "wd5~wd7: 9.16 mm/s",
            "  • 전압: 39.5V",
            "  • 전류: 255A",
            "  • 와이어 송급: 420 ipm",
            "",
            "wd8~wd9: 7.5 mm/s",
            "  • 전압: 36V",
            "  • 전류: 225A",
            "  • 와이어 송급: 375 ipm"
        ],
        right_title="위빙 패턴 (weavedata)",
        right_content=[
            "다양한 위빙 패턴 (40개 세트):",
            "",
            "weave1~4: 표준 패턴",
            "  • 폭: 2mm",
            "  • 높이: 2mm",
            "  • 빈도: 3Hz",
            "",
            "weave5~7: 확장 패턴",
            "  • 폭: 2mm",
            "  • 높이: 2mm",
            "  • 빈도: 3.5Hz",
            "",
            "weave8~9: 최소 패턴",
            "  • 위빙 없이 직선 용접",
            "",
            "속도 배열 (vWeld{40})",
            "  • 각 구간별 최적 속도 설정"
        ])

    # 12. UI 통신 구조
    slide = add_title_content_slide(prs, "외부 UI 통신 구조")
    add_two_column_content(slide,
        left_title="TASK8: 명령 인터페이스",
        left_content=[
            "PERS 변수 기반 통신:",
            "",
            "stCommand (명령 문자열)",
            "  • \"MoveJgJ\" - Joint 이동",
            "  • \"MovePgJ\" - Point → Joint",
            "  • \"MovePgL\" - Point → Linear",
            "  • \"WireCutR1/R2\" - 와이어 커팅",
            "",
            "stReact{3} (응답 배열)",
            "  • [\"Ready\", \"Ready\", \"Ready\"]",
            "  • [\"Ack\", \"Ack\", \"Ack\"]",
            "  • 3개 Task 응답 동기화",
            "",
            "데이터 교환:",
            "  • jRob1, jRob2, jGantry",
            "  • pRob1, pRob2",
            "  • vSync, zSync (속도/존)"
        ],
        right_title="I/O 신호 구성",
        right_content=[
            "CC_LINK (PLC 통신):",
            "  • co001~co010: 상태 신호",
            "  • AutoMode, PrgRun, MotorOn",
            "  • MovingRobot, MovingXY",
            "",
            "Local_IO (디지털 I/O):",
            "  • di09~16: 도어 센서 (4개 도어)",
            "  • di17~20: 충격 감지",
            "  • di21~29: 커터/노즐 상태",
            "  • do01~18: 도어/커터/냉각 제어",
            "",
            "Lincoln 용접기:",
            "  • soLn1/2StopProc (용접 중지)",
            "  • diLn1/2GasOK (가스 상태)",
            "  • diLn1/2WaterOK (냉각수)",
            "",
            "특수 신호:",
            "  • so_MoveG_PosHold (갠트리 위빙)",
            "  • pi01~18: 실시간 조정",
            "  • doR1/R2LeftSync (동기화)"
        ])

    # 13. 특징 및 제약사항
    slide = add_title_content_slide(prs, "특징 및 제약사항")
    add_two_column_content(slide,
        left_title="시스템 특징",
        left_content=[
            "✓ 3개 로봇 협조 제어",
            "",
            "✓ 소프트웨어 준동기화",
            "",
            "✓ 동적 TCP 변환",
            "",
            "✓ 위빙 중 갠트리 이동",
            "",
            "✓ 단계별 시퀀스 제어",
            "",
            "✓ 안전 리미트 체크"
        ],
        right_title="제약 사항",
        right_content=[
            "⚠ 하드웨어 동기화 불가",
            "  (MultiMove 2축 제한)",
            "",
            "⚠ 타이밍 민감도 높음",
            "  (신호 지연 영향)",
            "",
            "⚠ 복잡한 RAPID 로직",
            "  (유지보수 난이도)",
            "",
            "⚠ DeviceNet 미사용",
            "  (설정만 존재)",
            "",
            "⚠ 디버깅 어려움",
            "  (3개 Task 동시 분석)"
        ])

    # 14. 기술적 혁신
    slide = add_title_content_slide(prs, "기술적 혁신 포인트")
    add_bullet_points(slide, [
        "1. ABB MultiMove 2축 제약 극복",
        "   → RAPID 함수로 3-Task 준동기화 구현",
        "",
        "2. 동적 좌표 변환 알고리즘",
        "   → fnPoseToExtax + fnCoordToJoint 체인",
        "",
        "3. 신호 기반 모션 협조",
        "   → so_MoveG_PosHold로 상태 동기화",
        "",
        "4. 위빙 중 이동 제어",
        "   → WHILE 루프 + 느린 속도로 안정화",
        "",
        "5. 3단계 시퀀스 설계",
        "   → nWeldSequence 변수로 명확한 단계 구분"
    ])

    # 15. 시스템 장점
    slide = add_title_content_slide(prs, "시스템 장점")
    add_bullet_points(slide, [
        "생산성:",
        "  • 2개 로봇 동시 용접으로 사이클 타임 단축",
        "  • Gantry로 작업 영역 확장",
        "",
        "유연성:",
        "  • 소프트웨어 기반 제어로 수정 용이",
        "  • RAPID 파라미터 조정으로 최적화 가능",
        "",
        "안정성:",
        "  • 안전 리미트 체크 내장",
        "  • 단계별 시퀀스로 예측 가능한 동작",
        "",
        "확장성:",
        "  • 추가 위빙 패턴 구현 가능",
        "  • 새로운 용접 프로세스 적용 가능"
    ])

    # 16. 향후 개선 방향
    slide = add_title_content_slide(prs, "향후 개선 방향")
    add_bullet_points(slide, [
        "단기 개선:",
        "  1. DeviceNet 설정 정리 (사용 또는 제거)",
        "  2. RAPID 코드 리팩토링 (가독성 향상)",
        "  3. 디버깅 로그 추가 (문제 추적)",
        "",
        "중기 개선:",
        "  4. 자동 파라미터 튜닝 기능",
        "  5. 시뮬레이션 검증 도구 개발",
        "  6. 에러 핸들링 강화",
        "",
        "장기 개선:",
        "  7. RobotWare 7.x 업그레이드 검토",
        "  8. 외부 센서 통합 (비전, 레이저)",
        "  9. AI 기반 용접 품질 예측"
    ])

    # 17. 요약
    slide = add_title_content_slide(prs, "요약")
    add_bullet_points(slide, [
        "시스템 구성:",
        "  • 3개 로봇 (Robot1, Robot2, Gantry) 협조 제어",
        "  • ABB IRC5 + RobotWare 6.16 기반",
        "",
        "핵심 기술:",
        "  1. TCP 좌표 변환 (fnPoseToExtax, fnCoordToJoint)",
        "  2. 3-Task 준동기화 (WaitSyncTask, nWeldSequence)",
        "  3. 3단계 모션 시퀀스 (1→2→3)",
        "  4. 위빙 제어 (so_MoveG_PosHold + WHILE)",
        "",
        "특징:",
        "  • ABB MultiMove 2축 제약을 소프트웨어로 극복",
        "  • 안정적이고 확장 가능한 아키텍처",
        "  • 생산성과 유연성 동시 확보"
    ])

    # 18. 감사합니다
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (25, 25, 112), (65, 105, 225))

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "감사합니다"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(72)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Controller: 1200-526402 | Project: S25016"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(24)
    contact_para.font.color.rgb = RGBColor(200, 200, 255)
    contact_para.alignment = PP_ALIGN.CENTER

    return prs


def add_gradient_background(slide, rgb_start, rgb_end):
    """슬라이드에 그라디언트 배경 추가"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = RGBColor(*rgb_start)
    fill.gradient_stops[1].color.rgb = RGBColor(*rgb_end)


def add_title_content_slide(prs, title_text):
    """제목과 콘텐츠가 있는 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(25, 25, 112)

    # 제목 아래 구분선
    line = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1.2),
        Inches(15), Inches(0)
    )
    line.line.color.rgb = RGBColor(65, 105, 225)
    line.line.width = Pt(3)

    return slide


def add_bullet_points(slide, bullet_list):
    """슬라이드에 불릿 포인트 추가"""
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(14.4), Inches(6.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet_text in enumerate(bullet_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet_text
        p.level = 0 if not bullet_text.startswith("  ") else 1
        p.font.size = Pt(20) if p.level == 0 else Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)

        if bullet_text and not bullet_text.startswith(" "):
            p.font.bold = True


def add_two_column_content(slide, left_title, left_content, right_title, right_content):
    """2단 레이아웃으로 콘텐츠 추가"""
    # 왼쪽 컬럼
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(0.6))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.size = Pt(28)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = RGBColor(25, 25, 112)

    left_content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(6.5), Inches(5.5))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True

    for i, line in enumerate(left_content):
        if i == 0:
            p = left_content_frame.paragraphs[0]
        else:
            p = left_content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(6)
        if line and not line.startswith(" "):
            p.font.bold = True

    # 오른쪽 컬럼
    right_title_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(6.5), Inches(0.6))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.size = Pt(28)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = RGBColor(25, 25, 112)

    right_content_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.6), Inches(6.5), Inches(5.5))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True

    for i, line in enumerate(right_content):
        if i == 0:
            p = right_content_frame.paragraphs[0]
        else:
            p = right_content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(6)
        if line and not line.startswith(" "):
            p.font.bold = True


if __name__ == "__main__":
    print("ABB 로봇 시스템 기술 검토 PPT 생성 중...")
    prs = create_presentation()
    output_file = "ABB_Robot_System_Review_KO.pptx"
    prs.save(output_file)
    print(f"✓ PPT 파일이 생성되었습니다: {output_file}")
    print(f"  총 슬라이드: {len(prs.slides)}개")

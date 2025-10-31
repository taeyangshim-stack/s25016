#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Plan B 로봇 시스템 재구성 사양서 - PowerPoint 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """PowerPoint 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 1. 타이틀 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (139, 0, 0), (220, 20, 60))

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.text = "Plan B 로봇 시스템\n재구성 사양서"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(64)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "외부축 통합 방식\nController: 1200-526402 | 담당: 황동호 팀장, 유대영 차장\n기간: 2주 (2025-11-04 ~ 11-15)"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 200, 200)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 2. 프로젝트 개요
    slide = add_title_content_slide(prs, "프로젝트 개요")
    add_bullet_points(slide, [
        "프로젝트명: ABB 로봇 시스템 Plan B 구성",
        "",
        "목적:",
        "  • 현재 3-Task 준동기화 방식의 안정성 및 성능 개선",
        "  • Gantry를 Robot1의 외부축으로 통합하여 하드웨어 동기화 구현",
        "  • Robot2는 독립 제어하되 Task 간 협조로 동시 구동",
        "",
        "담당자:",
        "  • 황동호 팀장: 프로젝트 총괄, 시스템 설계, 하드웨어 설정",
        "  • 유대영 차장: RAPID 프로그래밍, 테스트 수행, 문서화",
        "",
        "일정:",
        "  • 시작: 2025-11-04 (예정)",
        "  • 종료: 2025-11-15 (예정)",
        "  • 총 기간: 2주 (10 영업일)"
    ])

    # 3. 시스템 구성 비교 - Plan A
    slide = add_title_content_slide(prs, "현재 구성 (Plan A)")
    add_bullet_points(slide, [
        "3개 메카니컬 유닛 독립 구성:",
        "",
        "TASK1 (T_ROB1): Robot1 제어",
        "  • IRB 1200-5/0.9",
        "  • Base Frame: [4665, -2217, 0]",
        "",
        "TASK2 (T_ROB2): Robot2 제어",
        "  • IRB 1200-5/0.9",
        "  • Base Frame: [4665, 2217, 0]",
        "",
        "TASK7 (Gantry_T): Gantry 제어",
        "  • 외부축 X, Y, Z, R (4축)",
        "  • fnPoseToExtax/fnCoordToJoint 좌표 변환 함수 사용",
        "",
        "동기화 방식: 소프트웨어 준동기화",
        "  • WaitSyncTask + nWeldSequence 변수",
        "  • so_MoveG_PosHold 신호 제어"
    ])

    # 4. 시스템 구성 비교 - Plan B
    slide = add_title_content_slide(prs, "새로운 구성 (Plan B)")
    add_bullet_points(slide, [
        "2개 메카니컬 유닛 + 외부축 통합:",
        "",
        "TASK1 (T_ROB1 + Gantry): Robot1 + 외부축 통합",
        "  • IRB 1200-5/0.9 (6축)",
        "  • 외부축 eax_a (X), eax_b (Y), eax_c (Z), eax_d (R)",
        "  • 하드웨어 동기화 (ABB MultiMove 지원)",
        "",
        "TASK2 (T_ROB2): Robot2 독립 제어",
        "  • IRB 1200-5/0.9 (6축)",
        "  • WaitSyncTask로 협조 동작",
        "",
        "장점:",
        "  ✓ Robot1과 Gantry 하드웨어 동기화",
        "  ✓ ABB MultiMove 2축 제한 준수",
        "  ✓ 좌표 변환 함수 불필요 (자동 처리)",
        "  ✓ RAPID 프로그래밍 단순화",
        "  ✓ 정밀도 향상 예상"
    ])

    # 5. 하드웨어 구성
    slide = add_title_content_slide(prs, "하드웨어 구성 상세")
    add_two_column_content(slide,
        left_title="메카니컬 유닛",
        left_content=[
            "ROB_1 (IRB 1200-5/0.9)",
            "  • 6축 로봇",
            "  • 외부축 4개 추가",
            "",
            "외부축 구성:",
            "  • eax_a: X축 (Linear)",
            "    Range: -5000 ~ +5000 mm",
            "",
            "  • eax_b: Y축 (Linear)",
            "    Range: -3000 ~ +3000 mm",
            "",
            "  • eax_c: Z축 (Linear)",
            "    Range: 0 ~ +2000 mm",
            "",
            "  • eax_d: R축 (Rotary)",
            "    Range: -180 ~ +180 deg",
            "",
            "ROB_2 (IRB 1200-5/0.9)",
            "  • 6축 로봇 (독립)"
        ],
        right_title="드라이브 할당",
        right_content=[
            "Drive Module 구성:",
            "",
            "ROB_1 (축 1-6)",
            "  • Drive Module 1-6",
            "  • 기존 로봇 모터 유지",
            "",
            "외부축 (eax_a/b/c/d)",
            "  • Drive Module 7-10",
            "  • 재구성 필요",
            "",
            "ROB_2 (축 1-6)",
            "  • Drive Module 11-16",
            "  • 기존 로봇 모터 유지",
            "",
            "확인 필요 사항:",
            "  ⚠ 드라이브 슬롯 가용성",
            "  ⚠ 전원 공급 용량",
            "  ⚠ 모터 사양 매칭"
        ])

    # 6. 외부축 운동학
    slide = add_title_content_slide(prs, "외부축 운동학 구성")
    add_bullet_points(slide, [
        "외부축 연결 체인:",
        "",
        "World Frame → eax_a (X) → eax_b (Y) → eax_c (Z) → eax_d (R) → ROB_1 Base",
        "",
        "각 축의 역할:",
        "  1. eax_a (X축): World Frame에 고정, 좌우 이동",
        "  2. eax_b (Y축): eax_a 위에 장착, 전후 이동",
        "  3. eax_c (Z축): eax_b 위에 장착, 상하 이동",
        "  4. eax_d (R축): eax_c 위에 장착, 회전 운동",
        "  5. ROB_1: eax_d 위에 장착, 6축 로봇 동작",
        "",
        "커플링 매트릭스:",
        "  • ARM.cfg에서 운동학 관계 정의",
        "  • MOC.cfg에서 메카니컬 유닛 연결",
        "  • 자동 좌표 변환 (Forward/Inverse Kinematics)"
    ])

    # 7. 소프트웨어 구성
    slide = add_title_content_slide(prs, "소프트웨어 구성")
    add_two_column_content(slide,
        left_title="RobotWare 옵션",
        left_content=[
            "필수 옵션:",
            "  • RobotWare 6.16 (유지)",
            "  • MultiMove Coordinated",
            "    - ROB_1 + 외부축",
            "    - ROB_2",
            "  • Arc Welding",
            "  • External Axis (4축)",
            "",
            "확인 필요:",
            "  ⚠ 현재 라이선스 확인",
            "  ⚠ 4축 외부축 지원 여부",
            "  ⚠ MultiMove 구성 변경",
            "",
            "추가 옵션 (유지):",
            "  • Collision Detection",
            "  • Path Recovery",
            "  • World Zones"
        ],
        right_title="시스템 파라미터 파일",
        right_content=[
            "수정 필요 파일:",
            "",
            "MOC.cfg",
            "  • ROB_1에 외부축 4개 추가",
            "  • Mechanical Unit 재정의",
            "",
            "ARM.cfg",
            "  • 외부축 운동학 정의",
            "  • 커플링 매트릭스 설정",
            "",
            "MOTOR.cfg",
            "  • 외부축 모터 파라미터",
            "  • 드라이브 매핑",
            "",
            "EIO.cfg",
            "  • I/O 신호 유지",
            "  • 일부 신호 재정의 가능",
            "",
            "PROC.cfg",
            "  • 용접 프로세스 유지"
        ])

    # 8. RAPID 프로그램 변경
    slide = add_title_content_slide(prs, "RAPID 프로그램 변경 사항")
    add_two_column_content(slide,
        left_title="변경 전 (Plan A)",
        left_content=[
            "TASK1 (T_ROB1):",
            "  MoveL Welds1{i}, vWeld{i},",
            "    z1, tWeld1, \\Weave:=weave;",
            "",
            "TASK7 (Gantry_T):",
            "  jTemp := fnCoordToJoint(",
            "    fnPoseToExtax(WeldsG{i}));",
            "  MoveExtJ jTemp, vWeld, fine;",
            "",
            "문제점:",
            "  • 좌표 변환 함수 복잡",
            "  • 3개 Task 동기화 어려움",
            "  • 타이밍 민감",
            "  • 디버깅 복잡"
        ],
        right_title="변경 후 (Plan B)",
        right_content=[
            "TASK1 (T_ROB1 + 외부축):",
            "  # extax 자동 동기화",
            "  MoveL Welds1{i}, vWeld{i},",
            "    z1, tWeld1, \\Weave:=weave;",
            "",
            "TASK7: 삭제",
            "  # Gantry가 ROB_1 외부축으로",
            "  # 통합되어 불필요",
            "",
            "개선 사항:",
            "  ✓ 좌표 변환 자동 처리",
            "  ✓ 2개 Task로 단순화",
            "  ✓ 하드웨어 동기화",
            "  ✓ 코드 가독성 향상",
            "",
            "robtarget 구조:",
            "  • trans: X, Y, Z",
            "  • rot: Quaternion",
            "  • extax: eax_a/b/c/d 자동"
        ])

    # 9. 좌표계 및 데이터 변환
    slide = add_title_content_slide(prs, "좌표계 및 데이터 변환")
    add_bullet_points(slide, [
        "Work Object 재정의:",
        "",
        "PERS wobjdata wobjGantry := [",
        "  FALSE,                    # robhold (외부축이 이동)",
        "  TRUE,                     # ufprog (사용자 좌표계 사용)",
        '  "",                       # ufmec',
        "  [[0,0,0],[1,0,0,0]],      # uframe (갠트리 기준 좌표계)",
        "  [[0,0,0],[1,0,0,0]]       # oframe (대상물 좌표계)",
        "];",
        "",
        "데이터 변환 작업:",
        "  1. WeldsG{} 배열을 Welds1{}.extax로 통합",
        "  2. 기존 Welds1{}.trans는 로봇 TCP 좌표 유지",
        "  3. Welds1{}.extax에 갠트리 위치 추가",
        "",
        "Tool 데이터:",
        "  • tWeld1, tWeld2 기존 유지",
        "  • TCP는 로봇 플랜지 기준 동일"
    ])

    # 10. 용접 파라미터 구성
    slide = add_title_content_slide(prs, "용접 파라미터 구성")
    add_two_column_content(slide,
        left_title="용접 조건 (유지)",
        left_content=[
            "welddata (40개 세트):",
            "",
            "wd1~4: 8.66 mm/s",
            "  [8.66, 0,",
            "   [5, 0, 40, 270, 0, 345, 0, 0, 0],",
            "   [0, 0, 0, 0, 0, 0, 0, 0, 0]]",
            "",
            "wd5~7: 9.16 mm/s",
            "  [9.16, 0,",
            "   [5, 0, 39.5, 255, 0, 420, 0, 0, 0],",
            "   [0, 0, 0, 0, 0, 0, 0, 0, 0]]",
            "",
            "wd8~9: 7.5 mm/s",
            "  [7.5, 0,",
            "   [5, 0, 36, 225, 0, 375, 0, 0, 0],",
            "   [0, 0, 0, 0, 0, 0, 0, 0, 0]]",
            "",
            "파라미터 의미:",
            "  • 속도 (mm/s)",
            "  • 전압 (V)",
            "  • 전류 (A)",
            "  • 와이어 송급 (ipm)"
        ],
        right_title="위빙 패턴 (유지)",
        right_content=[
            "weavedata (40개 세트):",
            "",
            "weave1~4: 표준 패턴",
            "  [1, 2, 2, 3, 0, 0, 0, ...]",
            "  • 타입: 1 (좌우 대칭)",
            "  • 폭: 2 mm",
            "  • 높이: 2 mm",
            "  • 빈도: 3 Hz",
            "",
            "weave5~7: 확장 패턴",
            "  [1, 2, 2, 3.5, 0, 0, 0, ...]",
            "  • 빈도: 3.5 Hz (증가)",
            "",
            "weave8~9: 직선 용접",
            "  [1, 0, 0, 0, 0, 0, 0, ...]",
            "  • 위빙 없음",
            "",
            "적용 방식:",
            "  • ArcL 명령어 \\Weave 옵션",
            "  • 외부축 이동 중에도 유지",
            "  • 하드웨어 동기화로 안정성↑"
        ])

    # 11. 용접 모션 시퀀스
    slide = add_title_content_slide(prs, "용접 모션 시퀀스 (Plan B)")
    add_bullet_points(slide, [
        "Phase 1: 로봇 진입 및 아크 시작",
        "  • ROB_1: 용접 시작점 이동 (외부축 정지)",
        "  • ROB_2: 동일하게 용접 시작점 이동",
        "  • WaitSyncTask로 동기화",
        "  • ArcL 명령으로 아크 점화",
        "",
        "Phase 2: 외부축 이동 + 위빙",
        "  • ROB_1: 외부축 자동 이동 + TCP 위빙",
        "    - 하드웨어 동기화로 정밀 제어",
        "    - MoveL \\Weave 한 번에 처리",
        "  • ROB_2: 협조 동작 (WaitSyncTask)",
        "",
        "Phase 3: 로봇 마무리 및 아크 소호",
        "  • ROB_1: 용접 종료점 도달, 외부축 정지",
        "  • ROB_2: 동일하게 종료점 도달",
        "  • 아크 소호",
        "",
        "개선 효과:",
        "  ✓ so_MoveG_PosHold 신호 불필요 (자동 동기화)",
        "  ✓ WHILE 루프 제거 (단순화)",
        "  ✓ 정밀도 향상"
    ])

    # 12. UI 통신 인터페이스
    slide = add_title_content_slide(prs, "UI 통신 인터페이스")
    add_two_column_content(slide,
        left_title="TASK8 명령 구조 (유지)",
        left_content=[
            "PERS 변수 기반 통신:",
            "",
            "stCommand (명령 문자열)",
            '  • "MoveJgJ" - Joint 이동',
            '  • "MovePgJ" - Point → Joint',
            '  • "MovePgL" - Point → Linear',
            '  • "WireCutR1/R2" - 와이어 커팅',
            "",
            "stReact{3} → stReact{2}",
            '  • ["Ready", "Ready"]',
            '  • ["Ack", "Ack"]',
            "  • 2개 Task로 변경 (TASK7 제거)",
            "",
            "데이터 교환 (수정):",
            "  • jRob1 (ROB_1 + 외부축)",
            "    - robax: 로봇 6축",
            "    - extax: 외부축 4축",
            "  • jRob2 (기존 유지)",
            "  • pRob1, pRob2 (기존 유지)",
            "  • vSync, zSync (유지)"
        ],
        right_title="I/O 신호 (유지)",
        right_content=[
            "CC_LINK (PLC 통신):",
            "  • co001_AutoMode",
            "  • co002_PrgRun",
            "  • co004_MotorOn",
            "  • co009_MovingRobot",
            "  • co010_MovingXY",
            "",
            "Local_IO (디지털 I/O):",
            "  • di09~16: 도어 센서",
            "  • di17~20: 충격 감지",
            "  • di21~29: 커터/노즐",
            "  • do01~18: 제어 신호",
            "",
            "Lincoln 용접기:",
            "  • soLn1/2StopProc",
            "  • diLn1/2GasOK",
            "  • diLn1/2WaterOK",
            "",
            "특수 신호 (변경):",
            "  • so_MoveG_PosHold: 제거 가능",
            "    (하드웨어 동기화로 불필요)",
            "  • pi01~18: 실시간 조정 (유지)",
            "  • doR1/R2LeftSync (유지)"
        ])

    # 13. UI 명령 처리 흐름
    slide = add_title_content_slide(prs, "UI 명령 처리 흐름 (Plan B)")
    add_bullet_points(slide, [
        "1. UI → RAPID 명령 전달:",
        "   stCommand := \"MovePgJ\"",
        "   pRob1 := [position with extax]",
        "   pRob2 := [position]",
        "",
        "2. TASK1 (ROB_1 + 외부축) 처리:",
        "   WaitUntil stCommand = \"MovePgJ\"",
        "   MoveL pRob1, vSync, zSync, tWeld1;  # 외부축 자동 이동",
        "   stReact{1} := \"Ack\"",
        "",
        "3. TASK2 (ROB_2) 처리:",
        "   WaitUntil stCommand = \"MovePgJ\"",
        "   MoveL pRob2, vSync, zSync, tWeld2;",
        "   stReact{2} := \"Ack\"",
        "",
        "4. UI 응답 확인:",
        "   WaitUntil stReact = [\"Ack\", \"Ack\"]",
        "   stCommand := \"\"",
        "   WaitUntil stReact = [\"Ready\", \"Ready\"]",
        "",
        "변경 사항:",
        "  • stReact{3} → stReact{2} (2개 Task)",
        "  • jGantry 변수 제거 (pRob1.extax로 통합)",
        "  • 좌표 변환 함수 호출 제거"
    ])

    # 14. 설정 작업 항목
    slide = add_title_content_slide(prs, "설정 작업 항목")
    add_bullet_points(slide, [
        "1. 시스템 백업 (Day 1)",
        "   • 전체 시스템 백업 생성",
        "   • RAPID 프로그램 백업",
        "   • 설정 파일 백업 (MOC.cfg, ARM.cfg, etc.)",
        "",
        "2. MOC.cfg 수정 (Day 2-3)",
        "   • ROB_1에 외부축 4개 추가",
        "   • Mechanical Unit 재정의",
        "   • ROB_2 유지",
        "",
        "3. ARM.cfg 작성 (Day 3)",
        "   • 외부축 운동학 정의",
        "   • 커플링 매트릭스 설정",
        "   • 각 축 DH 파라미터",
        "",
        "4. MOTOR.cfg 설정 (Day 3)",
        "   • 외부축 모터 파라미터",
        "   • 드라이브 매핑 (Drive 7-10)",
        "   • 속도/가감속 리미트",
        "",
        "5. RAPID 프로그램 수정 (Day 5-6)",
        "   • TASK7 삭제",
        "   • TASK1: robtarget.extax 추가",
        "   • TASK8: stReact{3} → {2} 변경"
    ])

    # 15. 캘리브레이션 절차
    slide = add_title_content_slide(prs, "캘리브레이션 절차")
    add_two_column_content(slide,
        left_title="축별 원점 설정 (Day 4)",
        left_content=[
            "1. X축 (eax_a) 원점:",
            "   • 기계적 원점 센서 확인",
            "   • 소프트 리미트 설정",
            "   • 홈 위치 정의",
            "",
            "2. Y축 (eax_b) 원점:",
            "   • X축 위 설치 기준 확인",
            "   • 원점 센서 확인",
            "   • 리미트 설정",
            "",
            "3. Z축 (eax_c) 원점:",
            "   • Y축 위 설치 기준 확인",
            "   • 최하단 위치 설정",
            "   • 안전 높이 확인",
            "",
            "4. R축 (eax_d) 원점:",
            "   • 0도 기준 위치 설정",
            "   • ROB_1 베이스 방향 확인",
            "   • 회전 범위 확인"
        ],
        right_title="동작 범위 확인 (Day 4)",
        right_content=[
            "소프트 리미트 설정:",
            "",
            "eax_a (X축):",
            "  • Min: -5000 mm",
            "  • Max: +5000 mm",
            "",
            "eax_b (Y축):",
            "  • Min: -3000 mm",
            "  • Max: +3000 mm",
            "",
            "eax_c (Z축):",
            "  • Min: 0 mm (안전 높이)",
            "  • Max: +2000 mm",
            "",
            "eax_d (R축):",
            "  • Min: -180 deg",
            "  • Max: +180 deg",
            "",
            "검증:",
            "  ✓ 하드 리미트 센서 동작",
            "  ✓ 비상정지 동작",
            "  ✓ 충돌 감지 활성화"
        ])

    # 16. 테스트 계획
    slide = add_title_content_slide(prs, "테스트 계획")
    add_bullet_points(slide, [
        "Phase 1: 외부축 단독 테스트 (Day 5)",
        "  • X, Y, Z, R 각 축 개별 동작 확인",
        "  • 리미트 및 홈 포지션 확인",
        "  • 속도 및 가감속 최적화",
        "  • 비상정지 및 안전 기능 확인",
        "",
        "Phase 2: ROB_1 + 외부축 통합 테스트 (Day 7)",
        "  • 로봇 단독 동작 (외부축 정지)",
        "  • 외부축 단독 동작 (로봇 정지)",
        "  • 동시 동작 (Linear, Circular 모션)",
        "  • 하드웨어 동기화 품질 확인",
        "",
        "Phase 3: ROB_1 + ROB_2 협조 테스트 (Day 8)",
        "  • WaitSyncTask 동기화 확인",
        "  • 용접 동작 재현",
        "  • 외부축 이동 중 위빙 테스트",
        "  • 타이밍 및 정밀도 측정",
        "",
        "Phase 4: 통합 시스템 테스트 (Day 9)",
        "  • 실제 용접 프로세스 재현",
        "  • UI 통신 확인",
        "  • 전체 시퀀스 동작 검증"
    ])

    # 17. 검증 항목
    slide = add_title_content_slide(prs, "검증 항목 및 합격 기준")
    add_two_column_content(slide,
        left_title="성능 검증",
        left_content=[
            "위치 정확도:",
            "  • 측정 방법: Laser tracker",
            "  • 합격 기준: ±0.5 mm",
            "  • 목표: ±0.3 mm",
            "",
            "반복 정밀도:",
            "  • 측정 방법: 동일 동작 10회",
            "  • 합격 기준: ±0.2 mm",
            "  • 목표: ±0.1 mm",
            "",
            "동기화 오차:",
            "  • 측정: ROB_1 vs ROB_2",
            "  • 합격 기준: <20 ms",
            "  • 목표: <10 ms",
            "",
            "사이클 타임:",
            "  • 측정: 실제 용접 프로세스",
            "  • 합격 기준: Plan A 대비 ±10%",
            "  • 목표: Plan A 대비 -5%"
        ],
        right_title="기능 검증",
        right_content=[
            "외부축 동작:",
            "  ✓ 4축 독립 제어 가능",
            "  ✓ 리미트 정상 동작",
            "  ✓ 홈 복귀 정상",
            "",
            "하드웨어 동기화:",
            "  ✓ ROB_1 + 외부축 동시 제어",
            "  ✓ 궤적 정확도 만족",
            "  ✓ 속도 프로파일 정상",
            "",
            "협조 동작:",
            "  ✓ WaitSyncTask 정상 동작",
            "  ✓ ROB_2 독립 제어 가능",
            "  ✓ 용접 품질 유지",
            "",
            "UI 통신:",
            "  ✓ 명령/응답 정상",
            "  ✓ stReact{2} 동작 확인",
            "  ✓ 데이터 교환 정상",
            "",
            "안전 기능:",
            "  ✓ 비상정지 동작",
            "  ✓ 충돌 감지 동작",
            "  ✓ 안전 속도 제한"
        ])

    # 18. 일정표 Week 1
    slide = add_title_content_slide(prs, "일정표 - Week 1 (설계 및 설정)")
    add_bullet_points(slide, [
        "Day 1 (11/04 월) - 킥오프 및 준비",
        "  • 킥오프 미팅 (2h) - 황동호, 유대영",
        "  • 시스템 전체 백업 (1h) - 유대영",
        "  • 외부축 파라미터 설계 초안 (3h) - 황동호",
        "  • MOC.cfg, ARM.cfg 초안 작성 (2h) - 유대영",
        "",
        "Day 2 (11/05 화) - 시스템 설계",
        "  • 드라이브 할당 계획 수립 (2h) - 황동호",
        "  • 외부축 모터 파라미터 확인 (2h) - 유대영",
        "  • RAPID 프로그램 수정 계획 (3h) - 유대영",
        "  • 캘리브레이션 절차 작성 (1h) - 황동호",
        "",
        "Day 3 (11/06 수) - 설정 파일 작성",
        "  • MOC.cfg, ARM.cfg 작성 및 검토 (4h) - 유대영",
        "  • MOTOR.cfg 설정 (2h) - 황동호",
        "  • 시스템 파라미터 적용 (2h) - 유대영",
        "",
        "Day 4 (11/07 목) - 캘리브레이션",
        "  • 외부축 원점 설정 X, Y, Z, R (3h) - 황동호, 유대영",
        "  • 소프트 리미트 설정 및 확인 (2h) - 유대영",
        "  • 커플링 매트릭스 설정 (2h) - 황동호",
        "  • Phase 1 테스트 준비 (1h) - 유대영",
        "",
        "Day 5 (11/08 금) - Phase 1 테스트",
        "  • 외부축 단독 테스트 (4h) - 황동호, 유대영",
        "  • 문제점 분석 및 수정 (2h) - 황동호, 유대영",
        "  • Week 1 진행 상황 보고 (1h) - 황동호",
        "  • RAPID 프로그램 수정 시작 (1h) - 유대영"
    ])

    # 19. 일정표 Week 2
    slide = add_title_content_slide(prs, "일정표 - Week 2 (테스트 및 검증)")
    add_bullet_points(slide, [
        "Day 6 (11/11 월) - RAPID 수정",
        "  • TASK1 통합 (외부축 추가) (4h) - 유대영",
        "  • fnPoseToExtax/fnCoordToJoint 제거 (1h) - 유대영",
        "  • robtarget.extax 데이터 변환 (2h) - 유대영",
        "  • 컴파일 및 오류 수정 (1h) - 유대영",
        "",
        "Day 7 (11/12 화) - Phase 2 테스트",
        "  • ROB_1 + 외부축 통합 테스트 (4h) - 황동호, 유대영",
        "  • Linear/Circular 모션 확인 (2h) - 유대영",
        "  • 동기화 품질 측정 (2h) - 황동호",
        "",
        "Day 8 (11/13 수) - Phase 3 테스트",
        "  • ROB_1 + ROB_2 협조 테스트 (4h) - 황동호, 유대영",
        "  • WaitSyncTask 동작 확인 (2h) - 유대영",
        "  • 외부축 이동 중 위빙 테스트 (2h) - 황동호, 유대영",
        "",
        "Day 9 (11/14 목) - Phase 4 테스트",
        "  • 통합 시스템 테스트 (3h) - 황동호, 유대영",
        "  • 실제 용접 프로세스 검증 (3h) - 황동호, 유대영",
        "  • 성능 측정 (정밀도, 사이클 타임) (2h) - 황동호",
        "",
        "Day 10 (11/15 금) - 최종 검증 및 보고",
        "  • 최종 검증 및 문서화 (3h) - 황동호, 유대영",
        "  • 문제점 목록 작성 (1h) - 유대영",
        "  • Plan A vs Plan B 비교 보고서 (2h) - 황동호",
        "  • 최종 결과 보고 (2h) - 황동호, 유대영"
    ])

    # 20. 마일스톤
    slide = add_title_content_slide(prs, "주요 마일스톤 및 산출물")
    add_two_column_content(slide,
        left_title="마일스톤",
        left_content=[
            "M1: 프로젝트 시작",
            "  • 날짜: 11/04 (Day 1)",
            "  • 산출물: 킥오프 의사록",
            "",
            "M2: 설정 파일 완료",
            "  • 날짜: 11/06 (Day 3)",
            "  • 산출물: MOC/ARM/MOTOR.cfg",
            "",
            "M3: Week 1 완료",
            "  • 날짜: 11/08 (Day 5)",
            "  • 산출물: 외부축 테스트 결과",
            "",
            "M4: RAPID 수정 완료",
            "  • 날짜: 11/12 (Day 7)",
            "  • 산출물: 통합 RAPID 프로그램",
            "",
            "M5: 협조 동작 확인",
            "  • 날짜: 11/13 (Day 8)",
            "  • 산출물: 협조 테스트 결과",
            "",
            "M6: 프로젝트 완료",
            "  • 날짜: 11/15 (Day 10)",
            "  • 산출물: 최종 보고서"
        ],
        right_title="산출물 목록",
        right_content=[
            "설계 문서:",
            "  □ 시스템 구성도",
            "  □ 외부축 운동학 설계",
            "  □ 드라이브 할당표",
            "",
            "설정 파일:",
            "  □ MOC.cfg (수정)",
            "  □ ARM.cfg (신규)",
            "  □ MOTOR.cfg (수정)",
            "",
            "RAPID 프로그램:",
            "  □ TASK1 (통합)",
            "  □ TASK2 (수정)",
            "  □ TASK8 (수정)",
            "",
            "절차서:",
            "  □ 캘리브레이션 절차서",
            "  □ 테스트 절차서",
            "",
            "보고서:",
            "  □ Week 1 진행 보고",
            "  □ Plan A vs Plan B 비교",
            "  □ 최종 결과 보고서",
            "  □ 문제점 및 개선 목록"
        ])

    # 21. 승인 및 확인 사항
    slide = add_title_content_slide(prs, "승인 및 확인 사항")
    add_bullet_points(slide, [
        "사전 확인 필요 항목:",
        "",
        "1. RobotWare 라이선스",
        "   □ 4축 외부축 지원 여부 확인",
        "   □ MultiMove 옵션 구성 확인",
        "   □ 추가 라이선스 구매 필요 여부",
        "",
        "2. 하드웨어 용량",
        "   □ 드라이브 모듈 가용 슬롯 확인 (4개 필요)",
        "   □ 전원 공급 용량 확인",
        "   □ 외부축 모터 사양 확인",
        "",
        "3. 안전 승인",
        "   □ 외부축 추가에 따른 안전 평가",
        "   □ 리스크 어세스먼트 재검토",
        "   □ 안전 펜스 및 라이트 커튼 확인",
        "",
        "4. 고객 승인",
        "   □ 시스템 설정 변경 사전 동의",
        "   □ 테스트 일정 조율",
        "   □ 다운타임 협의",
        "",
        "5. 백업 및 롤백",
        "   □ 전체 시스템 백업 확인",
        "   □ Plan A 복원 절차 준비",
        "   □ 비상 연락망 구성"
    ])

    # 22. Plan A vs Plan B 비교
    slide = add_title_content_slide(prs, "Plan A vs Plan B 상세 비교")
    add_two_column_content(slide,
        left_title="Plan A (현재)",
        left_content=[
            "구성:",
            "  • 3개 메카니컬 유닛",
            "  • 3개 RAPID Task",
            "",
            "장점:",
            "  ✓ 현재 안정적으로 작동",
            "  ✓ 설정 변경 불필요",
            "  ✓ 검증된 시스템",
            "",
            "단점:",
            "  ✗ 하드웨어 동기화 불가",
            "  ✗ 좌표 변환 함수 복잡",
            "  ✗ 타이밍 민감도 높음",
            "  ✗ 디버깅 어려움",
            "  ✗ 정밀도 제한",
            "",
            "유지보수:",
            "  • 복잡한 RAPID 로직",
            "  • 3개 Task 동시 분석 필요",
            "  • 신호 타이밍 이슈"
        ],
        right_title="Plan B (신규)",
        right_content=[
            "구성:",
            "  • 2개 메카니컬 유닛",
            "    (ROB_1 + 외부축 4개, ROB_2)",
            "  • 2개 RAPID Task",
            "",
            "장점:",
            "  ✓ 하드웨어 동기화",
            "  ✓ RAPID 단순화",
            "  ✓ 정밀도 향상 예상",
            "  ✓ ABB 표준 구성",
            "  ✓ 디버깅 용이",
            "",
            "단점:",
            "  ✗ 시스템 재구성 필요",
            "  ✗ 테스트 시간 소요",
            "  ✗ 외부축 설정 복잡도",
            "",
            "유지보수:",
            "  • 표준 ABB 구성",
            "  • 단순한 RAPID 로직",
            "  • 자동 좌표 변환"
        ])

    # 23. 기술적 고려사항
    slide = add_title_content_slide(prs, "기술적 고려사항")
    add_bullet_points(slide, [
        "1. 외부축 운동학 정확성",
        "   • DH 파라미터 정확한 측정 필요",
        "   • 커플링 매트릭스 검증 필요",
        "   • ABB RobotStudio 시뮬레이션 권장",
        "",
        "2. 드라이브 용량 및 성능",
        "   • 외부축 모터 토크/속도 사양 확인",
        "   • 드라이브 모듈 호환성 확인",
        "   • 과부하 보호 설정",
        "",
        "3. 하드웨어 동기화 품질",
        "   • 보간 주기(Interpolation cycle) 확인",
        "   • 궤적 정확도 측정 필요",
        "   • 속도 프로파일 최적화",
        "",
        "4. ROB_2 협조 동작",
        "   • WaitSyncTask 타이밍 재검증",
        "   • taskGroup12 동기화 품질 확인",
        "   • 독립 제어 vs 협조 제어 균형",
        "",
        "5. 용접 품질 유지",
        "   • 외부축 이동 중 위빙 안정성",
        "   • 아크 길이 제어 (CTSP)",
        "   • 용접 파라미터 재튜닝 가능성"
    ])

    # 24. 롤백 계획
    slide = add_title_content_slide(prs, "롤백 계획")
    add_two_column_content(slide,
        left_title="롤백 조건",
        left_content=[
            "다음 중 하나라도 발생 시:",
            "",
            "1. 성능 미달",
            "   • 위치 정확도 > ±0.5mm",
            "   • 반복 정밀도 > ±0.2mm",
            "   • 사이클 타임 +10% 초과",
            "",
            "2. 기능 오류",
            "   • 하드웨어 동기화 실패",
            "   • 외부축 제어 불능",
            "   • ROB_2 협조 동작 실패",
            "",
            "3. 안전 문제",
            "   • 충돌 감지 오동작",
            "   • 비상정지 실패",
            "   • 리미트 오류",
            "",
            "4. 일정 지연",
            "   • Day 8까지 Phase 2 미완료",
            "   • 심각한 기술 문제 발생",
            "",
            "5. 고객 요청",
            "   • 고객 불승인",
            "   • 긴급 생산 요청"
        ],
        right_title="롤백 절차",
        right_content=[
            "1단계: 즉시 중단",
            "   • 모든 테스트 중단",
            "   • 안전 상태 확인",
            "   • 시스템 정지",
            "",
            "2단계: 백업 복원",
            "   • Day 1 백업 파일 로드",
            "   • MOC.cfg, ARM.cfg 복원",
            "   • RAPID 프로그램 복원",
            "",
            "3단계: 시스템 재시작",
            "   • 제어기 재부팅",
            "   • 캘리브레이션 확인",
            "   • 홈 포지션 복귀",
            "",
            "4단계: Plan A 검증",
            "   • 3-Task 동작 확인",
            "   • 용접 프로세스 확인",
            "   • UI 통신 확인",
            "",
            "5단계: 문제 분석",
            "   • 실패 원인 분석",
            "   • 개선 방향 수립",
            "   • 재시도 계획 수립",
            "",
            "예상 소요 시간: 4-8 시간"
        ])

    # 25. 용어 정의
    slide = add_title_content_slide(prs, "용어 정의")
    add_two_column_content(slide,
        left_title="로봇 시스템 용어",
        left_content=[
            "Plan A:",
            "  • 현재 3-Task 준동기화 방식",
            "",
            "Plan B:",
            "  • Gantry를 ROB_1 외부축으로",
            "    통합하는 방식",
            "",
            "외부축 (External Axis):",
            "  • 로봇에 추가되는 선형/회전축",
            "  • eax_a/b/c/d (X/Y/Z/R)",
            "",
            "하드웨어 동기화:",
            "  • 제어기에서 운동학적으로",
            "    통합 제어",
            "  • Forward/Inverse Kinematics",
            "",
            "소프트웨어 준동기화:",
            "  • WaitSyncTask + 신호로 동기화",
            "  • 타이밍 기반 협조",
            "",
            "커플링 매트릭스:",
            "  • 외부축 간 운동학 관계 정의",
            "  • ARM.cfg에서 설정"
        ],
        right_title="RAPID 용어",
        right_content=[
            "robtarget:",
            "  • 로봇 목표 위치",
            "  • trans (X,Y,Z)",
            "  • rot (Quaternion)",
            "  • extax (외부축)",
            "",
            "jointtarget:",
            "  • 관절 각도 목표",
            "  • robax (로봇 6축)",
            "  • extax (외부축 4축)",
            "",
            "wobjdata:",
            "  • 작업 대상 좌표계",
            "  • robhold: 외부축 이동 여부",
            "  • ufprog: 사용자 좌표계",
            "",
            "tooldata:",
            "  • 툴 정의 (TCP)",
            "  • tWeld1, tWeld2",
            "",
            "WaitSyncTask:",
            "  • Task 간 동기화 명령",
            "  • taskGroup 지정",
            "",
            "ArcL:",
            "  • 아크 용접 선형 이동",
            "  • \\Weave 옵션"
        ])

    # 26. 참고 문서
    slide = add_title_content_slide(prs, "참고 문서 및 리소스")
    add_bullet_points(slide, [
        "ABB 공식 문서:",
        "  • RobotWare 6.16 Operating Manual",
        "  • Technical Reference Manual - RAPID Instructions",
        "  • Application Manual - MultiMove",
        "  • Application Manual - External Axes",
        "  • Product Manual - IRC5",
        "",
        "현재 시스템 백업:",
        "  • 1200-526402_BACKUP_2025-10-15-No_1",
        "  • RAPID 프로그램 (5,000+ 라인)",
        "  • 시스템 파라미터 파일",
        "",
        "관련 프로젝트 문서:",
        "  • ABB_Robot_System_Review_KO.pptx",
        "  • 로봇 시스템 기술 검토",
        "  • 용접 파라미터 설정",
        "  • UI 통신 인터페이스",
        "",
        "ABB 지원:",
        "  • ABB 코리아 기술 지원팀",
        "  • 온라인 포럼: ABB Robotics Forum",
        "  • RobotStudio 시뮬레이션"
    ])

    # 27. 요약
    slide = add_title_content_slide(prs, "프로젝트 요약")
    add_bullet_points(slide, [
        "프로젝트 목표:",
        "  • Gantry를 Robot1의 외부축으로 통합",
        "  • 하드웨어 동기화로 정밀도 향상",
        "  • RAPID 프로그래밍 단순화",
        "",
        "주요 작업:",
        "  • MOC.cfg, ARM.cfg, MOTOR.cfg 설정",
        "  • 외부축 4개 캘리브레이션",
        "  • RAPID 프로그램 수정 (TASK7 제거, TASK1 통합)",
        "  • 4단계 테스트 (Phase 1~4)",
        "",
        "담당 및 일정:",
        "  • 황동호 팀장: 시스템 설계, 하드웨어 설정",
        "  • 유대영 차장: RAPID 프로그래밍, 테스트",
        "  • 2주 (2025-11-04 ~ 11-15)",
        "",
        "기대 효과:",
        "  • 위치 정확도 향상 (±0.5mm → ±0.3mm 목표)",
        "  • 유지보수성 개선 (RAPID 단순화)",
        "  • ABB 표준 구성 준수",
        "",
        "안전장치:",
        "  • 시스템 백업 및 롤백 계획",
        "  • 단계별 테스트 및 검증"
    ])

    # 28. 승인란
    slide = add_title_content_slide(prs, "승인 및 서명")
    add_bullet_points(slide, [
        "문서 정보:",
        "  • 문서명: Plan B 로봇 시스템 재구성 사양서",
        "  • 버전: 1.0",
        "  • 작성일: 2025-10-31",
        "",
        "작성자:",
        "  • 이름: _____________________",
        "  • 서명: _____________________",
        "  • 날짜: _____________________",
        "",
        "검토자:",
        "  • 황동호 팀장",
        "  • 서명: _____________________",
        "  • 날짜: _____________________",
        "",
        "  • 유대영 차장",
        "  • 서명: _____________________",
        "  • 날짜: _____________________",
        "",
        "승인자:",
        "  • 이름: _____________________",
        "  • 직책: _____________________",
        "  • 서명: _____________________",
        "  • 날짜: _____________________"
    ])

    # 29. 감사합니다
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, (139, 0, 0), (220, 20, 60))

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "감사합니다"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(72)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Plan B 프로젝트\n담당: 황동호 팀장, 유대영 차장\n2025-11-04 ~ 11-15 (2주)"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(24)
    contact_para.font.color.rgb = RGBColor(255, 200, 200)
    contact_para.alignment = PP_ALIGN.CENTER

    return prs


def add_gradient_background(slide, rgb_start, rgb_end):
    """슬라이드에 그라디언트 배경 추가"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = RGBColor(*rgb_start)
    fill.gradient_stops[1].color.rgb = RGBColor(*rgb_end)


def add_title_content_slide(prs, title_text):
    """제목과 콘텐츠가 있는 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 245, 245)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(139, 0, 0)

    # 제목 아래 구분선
    line = slide.shapes.add_shape(
        1,  # Line
        Inches(0.5), Inches(1.2),
        Inches(15), Inches(0)
    )
    line.line.color.rgb = RGBColor(220, 20, 60)
    line.line.width = Pt(3)

    return slide


def add_bullet_points(slide, bullet_list):
    """슬라이드에 불릿 포인트 추가"""
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(14.4), Inches(6.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet_text in enumerate(bullet_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet_text
        p.level = 0 if not bullet_text.startswith("  ") else 1
        p.font.size = Pt(18) if p.level == 0 else Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(6)

        if bullet_text and not bullet_text.startswith(" "):
            p.font.bold = True


def add_two_column_content(slide, left_title, left_content, right_title, right_content):
    """2단 레이아웃으로 콘텐츠 추가"""
    # 왼쪽 컬럼
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(0.6))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.size = Pt(26)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = RGBColor(139, 0, 0)

    left_content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(6.5), Inches(5.5))
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True

    for i, line in enumerate(left_content):
        if i == 0:
            p = left_content_frame.paragraphs[0]
        else:
            p = left_content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(4)
        if line and not line.startswith(" "):
            p.font.bold = True

    # 오른쪽 컬럼
    right_title_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(6.5), Inches(0.6))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.size = Pt(26)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = RGBColor(139, 0, 0)

    right_content_box = slide.shapes.add_textbox(Inches(8.2), Inches(2.6), Inches(6.5), Inches(5.5))
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True

    for i, line in enumerate(right_content):
        if i == 0:
            p = right_content_frame.paragraphs[0]
        else:
            p = right_content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(4)
        if line and not line.startswith(" "):
            p.font.bold = True


if __name__ == "__main__":
    print("Plan B 로봇 시스템 재구성 사양서 PPT 생성 중...")
    prs = create_presentation()
    output_file = "Plan_B_Robot_Reconfiguration_Spec.pptx"
    prs.save(output_file)
    print(f"✓ PPT 파일이 생성되었습니다: {output_file}")
    print(f"  총 슬라이드: {len(prs.slides)}개")
